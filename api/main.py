"""Минимальный FastAPI-бэкенд приложения участника.

Запускается как sidecar в pod'е приложения, nginx-фронт проксирует
`/api/*` сюда. Зачем существует: даёт фронту публичный endpoint для
вызовов, которые нельзя сделать из браузера напрямую — например,
LLM/VLM запросы (нужен серверный ключ к gpt2giga proxy).

Эндпоинты:
  GET  /api/healthz — health check.
  POST /api/vlm     — multipart image + prompt → текст распознавания
                      (через gpt2giga → GigaChat-Vision).
  POST /api/llm     — JSON {prompt | messages, system?, max_tokens?}
                      → {"text": "..."} текстовый ответ модели
                      (через gpt2giga → GigaChat-2-Max).
  POST /api/pptx    — JSON {title, subtitle?, slides:[{heading, bullets|body}]}
                      → файл .pptx attachment'ом.
  POST /api/image   — JSON {prompt, model?} → {"image": "data:image/png;base64,..."}
                      Генерация картинки через кластерный LiteLLM
                      (Gemini image-моделями).

Расширение: если участник просит ручку, которой здесь нет —
добавь её прямо в этот файл, deploy.sh пересоберёт image.
"""

import base64
import logging
import os
from io import BytesIO
from typing import Literal

import httpx
from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.responses import StreamingResponse
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pydantic import BaseModel, Field

VLM_URL = os.environ.get(
    "VLM_URL", "http://gpt2giga.ranepa-tools.svc.cluster.local:8090"
)
VLM_MODEL = os.environ.get("VLM_MODEL", "GigaChat-2-Max")
MAX_IMAGE_BYTES = 8 * 1024 * 1024  # 8 MB
ALLOWED_MEDIA = {"image/jpeg", "image/png", "image/gif", "image/webp"}

# LiteLLM gateway (кластерный) — для генерации картинок Gemini-моделями.
# Ключ и base-url приходят через env (k8s secret litellm-creds).
LITELLM_URL = os.environ.get(
    "LITELLM_URL", "http://litellm.openclaw.svc.cluster.local:4000"
)
LITELLM_KEY = os.environ.get("LITELLM_KEY", "")
# primary,fallback — пробуем по очереди (3.1-preview иногда нет в регионе).
IMAGE_MODELS = os.environ.get(
    "IMAGE_MODELS", "gemini-2.5-flash-image,gemini-3.1-flash-image-preview"
).split(",")

logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")
log = logging.getLogger("api")

app = FastAPI(title="ranepa-app api")


@app.get("/api/healthz")
async def healthz() -> dict[str, bool]:
    return {"ok": True}


@app.post("/api/vlm")
async def vlm(
    image: UploadFile = File(...),
    prompt: str = Form("Опиши изображение детально на русском."),
) -> dict[str, str]:
    raw = await image.read()
    if len(raw) > MAX_IMAGE_BYTES:
        raise HTTPException(413, "image too large (max 8 MB)")

    media_type = image.content_type or "image/jpeg"
    if media_type not in ALLOWED_MEDIA:
        raise HTTPException(400, f"unsupported media type: {media_type}")

    data = base64.b64encode(raw).decode("ascii")

    payload = {
        "model": VLM_MODEL,
        "max_tokens": 2048,
        "messages": [
            {
                "role": "user",
                "content": [
                    {
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": media_type,
                            "data": data,
                        },
                    },
                    {"type": "text", "text": prompt},
                ],
            }
        ],
    }

    async with httpx.AsyncClient(timeout=60) as client:
        r = await client.post(f"{VLM_URL}/v1/messages", json=payload)

    if r.status_code != 200:
        log.warning("vlm upstream %s: %s", r.status_code, r.text[:300])
        raise HTTPException(r.status_code, f"vlm upstream: {r.text[:500]}")

    body = r.json()
    text = "".join(
        b.get("text", "") for b in body.get("content", []) if b.get("type") == "text"
    )
    return {"text": text}


class Message(BaseModel):
    role: Literal["user", "assistant"]
    content: str


class LLMRequest(BaseModel):
    # Один из двух: либо короткий `prompt` (single-turn), либо
    # `messages` (multi-turn, фронт сам ведёт историю).
    prompt: str | None = None
    messages: list[Message] | None = None
    system: str | None = None
    max_tokens: int = Field(default=2048, ge=1, le=8192)


class LLMResponse(BaseModel):
    text: str


@app.post("/api/llm")
async def llm(req: LLMRequest) -> LLMResponse:
    if req.messages and req.prompt:
        raise HTTPException(400, "use either `prompt` or `messages`, not both")
    if not req.messages and not req.prompt:
        raise HTTPException(400, "either `prompt` or `messages` is required")

    messages = (
        [{"role": m.role, "content": m.content} for m in req.messages]
        if req.messages
        else [{"role": "user", "content": req.prompt}]
    )

    payload: dict[str, object] = {
        "model": VLM_MODEL,
        "max_tokens": req.max_tokens,
        "messages": messages,
    }
    if req.system:
        payload["system"] = req.system

    async with httpx.AsyncClient(timeout=60) as client:
        r = await client.post(f"{VLM_URL}/v1/messages", json=payload)

    if r.status_code != 200:
        log.warning("llm upstream %s: %s", r.status_code, r.text[:300])
        raise HTTPException(r.status_code, f"llm upstream: {r.text[:500]}")

    body = r.json()
    text = "".join(
        b.get("text", "") for b in body.get("content", []) if b.get("type") == "text"
    )
    return LLMResponse(text=text)


class PPTXSlide(BaseModel):
    heading: str
    bullets: list[str] = []
    body: str | None = None


class PPTXRequest(BaseModel):
    title: str
    subtitle: str | None = None
    slides: list[PPTXSlide] = Field(default_factory=list, max_length=30)
    filename: str = "presentation.pptx"


# Цвета — повторяем визуальный язык нашей лекторской колоды,
# чтобы участник сразу получал что-то стилистически похожее.
_BG = RGBColor(0x1F, 0x29, 0x3B)        # navy фон
_ACCENT = RGBColor(0x31, 0xBD, 0xD9)    # циан акцент
_TEXT = RGBColor(0xFF, 0xFF, 0xFF)      # белый
_MUTED = RGBColor(0xA9, 0xB8, 0xCE)     # приглушённо-голубой


def _set_slide_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _style_run(run, *, size: int, color: RGBColor, bold: bool = False) -> None:
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Arial"


@app.post("/api/pptx")
async def pptx(req: PPTXRequest) -> StreamingResponse:
    """Собирает .pptx из переданной структуры и отдаёт как download.

    По умолчанию — 5 слайдов на приложение: о приложении, экраны,
    как пользоваться, технологии, контакты. Но API принимает любое
    количество (до 30) — структура произвольная.
    """
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)

    # Title slide — blank layout, рисуем руками.
    title_slide = p.slides.add_slide(p.slide_layouts[6])  # blank
    _set_slide_bg(title_slide, _BG)

    title_box = title_slide.shapes.add_textbox(
        Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.8)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.text = req.title
    _style_run(tf.paragraphs[0].runs[0], size=54, color=_ACCENT, bold=True)

    if req.subtitle:
        sub_box = title_slide.shapes.add_textbox(
            Inches(0.8), Inches(4.6), Inches(11.7), Inches(0.9)
        )
        sub_tf = sub_box.text_frame
        sub_tf.word_wrap = True
        sub_tf.text = req.subtitle
        _style_run(sub_tf.paragraphs[0].runs[0], size=24, color=_MUTED)

    # Content slides.
    for slide_spec in req.slides:
        s = p.slides.add_slide(p.slide_layouts[6])
        _set_slide_bg(s, _BG)

        # Заголовок.
        title_box = s.shapes.add_textbox(
            Inches(0.7), Inches(0.5), Inches(12), Inches(1.1)
        )
        ttf = title_box.text_frame
        ttf.word_wrap = True
        ttf.text = slide_spec.heading
        _style_run(ttf.paragraphs[0].runs[0], size=36, color=_ACCENT, bold=True)

        # Тело.
        body_box = s.shapes.add_textbox(
            Inches(0.9), Inches(1.9), Inches(11.5), Inches(5.0)
        )
        btf = body_box.text_frame
        btf.word_wrap = True

        if slide_spec.bullets:
            for i, b in enumerate(slide_spec.bullets):
                para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
                para.text = f"•  {b}"
                para.space_after = Pt(12)
                _style_run(para.runs[0], size=22, color=_TEXT)
        elif slide_spec.body:
            btf.text = slide_spec.body
            _style_run(btf.paragraphs[0].runs[0], size=22, color=_TEXT)

    buf = BytesIO()
    p.save(buf)
    buf.seek(0)

    safe_name = req.filename if req.filename.endswith(".pptx") else f"{req.filename}.pptx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{safe_name}"'},
    )


class ImageRequest(BaseModel):
    prompt: str = Field(min_length=1, max_length=2000)
    model: str | None = None  # переопределить модель (по умолчанию — IMAGE_MODELS)


class ImageResponse(BaseModel):
    image: str  # data:image/png;base64,...


@app.post("/api/image")
async def image(req: ImageRequest) -> ImageResponse:
    """Генерация картинки по текстовому описанию через кластерный LiteLLM.

    Возвращает data-URL (можно сразу вставить в <img src>). Перебирает
    модели из IMAGE_MODELS по очереди — если первая недоступна в регионе
    Vertex, пробует следующую.
    """
    if not LITELLM_KEY:
        raise HTTPException(503, "image generation not configured (no LITELLM_KEY)")

    models = [req.model] if req.model else IMAGE_MODELS
    last_err = "no models tried"

    async with httpx.AsyncClient(timeout=120) as client:
        for model in models:
            try:
                r = await client.post(
                    f"{LITELLM_URL}/v1/images/generations",
                    headers={"Authorization": f"Bearer {LITELLM_KEY}"},
                    json={"model": model, "prompt": req.prompt},
                )
            except httpx.HTTPError as e:
                last_err = f"{model}: {e}"
                continue
            if r.status_code != 200:
                last_err = f"{model}: HTTP {r.status_code} {r.text[:200]}"
                continue
            body = r.json()
            data = (body.get("data") or [{}])[0]
            b64 = data.get("b64_json")
            if b64:
                return ImageResponse(image=f"data:image/png;base64,{b64}")
            url = data.get("url")
            if url:
                # Некоторые модели отдают URL вместо base64 — скачаем и
                # перекодируем, чтобы клиент всегда получал data-URL.
                img = await client.get(url)
                if img.status_code == 200:
                    enc = base64.b64encode(img.content).decode("ascii")
                    return ImageResponse(image=f"data:image/png;base64,{enc}")
            last_err = f"{model}: no image in response"

    log.warning("image gen failed: %s", last_err)
    raise HTTPException(502, f"image gen failed: {last_err}")
